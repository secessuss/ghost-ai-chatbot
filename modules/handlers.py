import asyncio
import re
import io
import logging
import fitz
import docx
import pandas as pd
import pptx
from telebot import types
from telebot.asyncio_helper import ApiTelegramException
from .bot_setup import bot, context_manager 
from .llm_text import generate_response_stream, _summarize_text_async
from .llm_vision import generate_response_from_image_stream
from .voice_handler import process_voice_message
from .web_handler import extract_content_from_url
from .utils import send_or_edit_message


TELEGRAM_MAX_MESSAGE_LENGTH = 4096
SUMMARY_THRESHOLD = 3800
URL_REGEX = r'(https?://\S+)'
stop_requests = {}

def escape_html(text: str) -> str:
    if not isinstance(text, str):
        return ""
    return text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')

def convert_markdown_to_html_safely(text: str) -> str:
    text = text.strip()
    text = re.sub(r'^\s*[\*\-]\s+', '', text, flags=re.MULTILINE)
    
    token_regex = re.compile(
        r'(```(?:[a-zA-Z]+\n)?(?:.|\n)*?```)|'  # Code blocks
        r'(\*\*.*?\*\*|__.*?__)|'               # Bold
        r'(\*.*?\*|_.*?_)|'                     # Italic
        r'(`.*?`)'                              # Inline code
    )
    
    html_parts = []
    last_end = 0
    
    for match in token_regex.finditer(text):
        start, end = match.span()
        html_parts.append(escape_html(text[last_end:start]))
        last_end = end
        token = match.group(0)
        
        if match.group(1):  # Code block
            inner_content = re.sub(r'^```[a-zA-Z]*\n?|```$', '', token)
            html_parts.append(f'<pre><code>{escape_html(inner_content.strip())}</code></pre>')
        elif match.group(2):  # Bold
            html_parts.append(f'<b>{escape_html(token.strip("*_"))}</b>')
        elif match.group(3):  # Italic
            html_parts.append(f'<i>{escape_html(token.strip("*_"))}</i>')
        elif match.group(4):  # Inline code
            html_parts.append(f'<code>{escape_html(token.strip("`"))}</code>')
    
    html_parts.append(escape_html(text[last_end:]))
    return "".join(html_parts)

def split_message(html_text: str, limit: int = TELEGRAM_MAX_MESSAGE_LENGTH) -> list[str]:
    if len(html_text) <= limit:
        return [html_text]
    
    final_chunks = []
    parts = re.split(r'(<pre><code>.*?</code></pre>)', html_text, flags=re.DOTALL)

    for part in parts:
        if not part:
            continue
            
        if part.startswith('<pre'):
            if len(part) > limit:
                final_chunks.extend([part[i:i + limit] for i in range(0, len(part), limit)])
            else:
                final_chunks.append(part)
            continue

        chunks = []
        current_chunk = ""
        open_tags = []
        sub_parts = re.split(r'(<[^>]+>)', part)
        
        for sub_part in sub_parts:
            if not sub_part:
                continue
                
            if len(current_chunk) + len(sub_part) > limit:
                for open_tag in reversed(open_tags):
                    tag_name = open_tag.strip('<>').split()[0]
                    current_chunk += f"</{tag_name}>"
                chunks.append(current_chunk)
                current_chunk = "".join(open_tags)
            
            current_chunk += sub_part

            if sub_part.startswith('<') and not sub_part.startswith('</'):
                if not sub_part.endswith('/>'):
                    open_tags.append(sub_part)
            elif sub_part.startswith('</'):
                tag_name_to_close = sub_part.strip('</> ')
                for i in range(len(open_tags) - 1, -1, -1):
                    open_tag = open_tags[i]
                    open_tag_name = open_tag.strip('<>').split()[0]
                    if open_tag_name == tag_name_to_close:
                        open_tags.pop(i)
                        break

        if current_chunk:
            chunks.append(current_chunk)
        final_chunks.extend(chunks)

    return final_chunks if final_chunks else [html_text]

async def stream_response(chat_id: int, placeholder_msg: types.Message, response_generator):
    full_response = ""
    last_edit_time = 0
    typing_cursor = "▌"
    stop_requests[chat_id] = False
    
    try:
        async for chunk in response_generator:
            if stop_requests.get(chat_id):
                full_response += "\n\n_Proses dihentikan oleh pengguna._"
                break
                
            if isinstance(chunk, dict):
                if 'event' in chunk:
                    status_map = {
                        'RESEARCH_START': '🤖🔍 Mencari informasi...',
                        'GENERATION_START': '🤖💬 Memberikan respon...',
                        'EXTRACTING_DESCRIPTION': '🤖🖼️ Memproses permintaan gambar...',
                        'GENERATING_PROMPT': "🤖⏳ Memproses deskripsi...",
                        'GENERATING_IMAGE': '🤖🎨 Membuat gambar...'
                    }
                    status_text = status_map.get(chunk['event'])
                    if status_text:
                        placeholder_msg = await send_or_edit_message(
                            chat_id, status_text, placeholder_msg, parse_mode='HTML'
                        )
                    
                    if chunk['event'] == 'RESEARCH_QUERY':
                        query_text = escape_html(chunk.get('data', '...'))
                        placeholder_msg = await send_or_edit_message(
                            chat_id, 
                            f"🤖🔎 Mencari: <i>\"{query_text}\"</i>", 
                            placeholder_msg, 
                            parse_mode='HTML'
                        )
                    continue
                    
                elif chunk['type'] == 'image':
                    await bot.delete_message(chat_id, placeholder_msg.message_id)
                    await bot.send_photo(
                        chat_id, 
                        chunk['data'], 
                        caption=f"🖼️ <code>{escape_html(chunk['caption'])}</code>", 
                        parse_mode="HTML"
                    )
                    return
                    
                elif chunk['type'] == 'stream_end':
                    full_response = chunk['full_text']
                    break
                    
                elif chunk['type'] == 'text':
                    data = chunk.get('data', '')
                    if isinstance(data, str):
                        full_response += data
                    else:
                        full_response = str(data)
                        break
            
            current_time = asyncio.get_event_loop().time()
            if current_time - last_edit_time > 1.2 and full_response:
                markup = types.InlineKeyboardMarkup().add(
                    types.InlineKeyboardButton("⏹️ Hentikan", callback_data=f"stop_{chat_id}")
                )
                placeholder_msg = await send_or_edit_message(
                    chat_id, 
                    full_response + typing_cursor, 
                    placeholder_msg, 
                    parse_mode='Markdown', 
                    reply_markup=markup
                )
                last_edit_time = current_time

        final_text_to_send = full_response
        if len(full_response) > SUMMARY_THRESHOLD:
            placeholder_msg = await send_or_edit_message(
                chat_id, "🤖🔍 Menganalisis...", placeholder_msg
            )
            final_text_to_send = await _summarize_text_async(full_response)

        final_html = convert_markdown_to_html_safely(final_text_to_send.strip())
        message_chunks = split_message(final_html)
        
        current_placeholder = placeholder_msg
        for i, text_chunk in enumerate(message_chunks):
            if i == 0:
                current_placeholder = await send_or_edit_message(
                    chat_id, text_chunk, current_placeholder, parse_mode='HTML'
                )
            else:
                await bot.send_message(chat_id, text_chunk, parse_mode='HTML')
                
    except Exception as e:
        logging.error(f"Error tak terduga saat menangani respons stream: {e}", exc_info=True)
        await send_or_edit_message(
            chat_id, "❗🤖 Terjadi kesalahan saat menampilkan respons.", placeholder_msg
        )
    finally:
        if chat_id in stop_requests:
            del stop_requests[chat_id]

def _get_main_menu_markup():
    markup = types.InlineKeyboardMarkup(row_width=2)
    buttons = [
        types.InlineKeyboardButton("⚙️ Kelola", callback_data="menu_sesi"),
        types.InlineKeyboardButton("🗑️ Hapus", callback_data="menu_reset"),
        types.InlineKeyboardButton("📖 Bantuan", callback_data="menu_help"),
        types.InlineKeyboardButton("❎ Tutup", callback_data="menu_tutup")
    ]
    markup.add(*buttons)
    return markup

def _get_session_menu_markup():
    markup = types.InlineKeyboardMarkup(row_width=2)
    buttons = [
        types.InlineKeyboardButton("🟢 Sesi Aktif", callback_data="sesi_aktif"),
        types.InlineKeyboardButton("🔴 Akhiri Sesi", callback_data="sesi_akhir"),
        types.InlineKeyboardButton("⬅️ Kembali", callback_data="menu_kembali")
    ]
    markup.add(*buttons)
    return markup

def _get_back_to_main_menu_markup():
    markup = types.InlineKeyboardMarkup()
    markup.add(types.InlineKeyboardButton("⬅️ Kembali ke Menu", callback_data="menu_kembali"))
    return markup

def _get_back_to_session_menu_markup():
    markup = types.InlineKeyboardMarkup()
    markup.add(types.InlineKeyboardButton("⬅️ Kembali", callback_data="menu_sesi"))
    return markup

@bot.message_handler(commands=['start'])
async def handle_start_command(message):
    await context_manager.reset_context(message.from_user.id)
    await bot.send_message(message.chat.id, "<b>🤖👋 Halo! Saya GHOST.</b>", parse_mode="HTML")

@bot.message_handler(commands=['menu'])
async def handle_menu_command(message):
    markup = _get_main_menu_markup()
    await bot.send_message(
        message.chat.id, 
        "<b>⚙️ Menu GHOST</b>\n\nPilih fitur di bawah:", 
        reply_markup=markup, 
        parse_mode="HTML"
    )

async def process_user_query(message, query_text):
    if not query_text or len(query_text.strip()) < 2:
        await bot.reply_to(message, "🤖 Pertanyaan kurang jelas.")
        return
        
    placeholder_msg = await bot.reply_to(message, "🤖💭 Berpikir...", parse_mode='HTML')
    response_generator = generate_response_stream(message.from_user.id, query_text)
    await stream_response(message.chat.id, placeholder_msg, response_generator)

@bot.message_handler(content_types=['text'])
async def handle_text(message):
    if message.text.startswith('/'):
        return
        
    urls = re.findall(URL_REGEX, message.text)
    if urls and len(message.text.strip()) == len(urls[0]):
        await handle_url_message(message, urls[0])
    else:
        await process_user_query(message, message.text)

@bot.message_handler(content_types=['voice', 'photo', 'document'])
async def handle_media(message):
    if message.content_type == 'voice':
        placeholder = await bot.reply_to(message, "🤖🎙️ Memproses pesan suara...", parse_mode='HTML')
        transcribed_text = await process_voice_message(message)
        
        if transcribed_text:
            await send_or_edit_message(
                message.chat.id, 
                f"<b>📝 Transkripsi:</b>\n<i>\"{escape_html(transcribed_text)}\"</i>", 
                placeholder, 
                parse_mode='HTML'
            )
            await process_user_query(message, transcribed_text)
        else:
            await send_or_edit_message(
                message.chat.id, "❌ Gagal memproses pesan suara.", placeholder
            )
            
    elif message.content_type == 'photo':
        prompt = message.caption or "Jelaskan gambar ini secara detail."
        placeholder_msg = await bot.reply_to(message, "🤖🖼️ Menganalisis gambar...", parse_mode='HTML')
        
        file_info = await bot.get_file(message.photo[-1].file_id)
        image_bytes = await bot.download_file(file_info.file_path)
        response_generator = generate_response_from_image_stream(
            message.from_user.id, prompt, image_bytes
        )
        await stream_response(message.chat.id, placeholder_msg, response_generator)
        
    elif message.content_type == 'document':
        await handle_document(message)

async def handle_document(message):
    if message.document.file_size > 5 * 1024 * 1024:
        await bot.reply_to(message, "📁 File terlalu besar (maks 5 MB).")
        return
        
    placeholder = await bot.reply_to(
        message, 
        f"📄 <i>Memproses <code>{escape_html(message.document.file_name)}</code>...</i>", 
        parse_mode='HTML'
    )
    
    try:
        file_info = await bot.get_file(message.document.file_id)
        downloaded_file = await bot.download_file(file_info.file_path)
        file_name_lower = message.document.file_name.lower()
        
        def process_file_sync():
            if file_name_lower.endswith('.pdf'):
                with fitz.open(stream=io.BytesIO(downloaded_file), filetype="pdf") as doc:
                    return "".join(page.get_text() for page in doc)
                    
            elif file_name_lower.endswith('.docx'):
                doc = docx.Document(io.BytesIO(downloaded_file))
                return "\n".join(para.text for para in doc.paragraphs)
                
            elif file_name_lower.endswith('.pptx'):
                prs = pptx.Presentation(io.BytesIO(downloaded_file))
                return "\n".join(
                    shape.text for slide in prs.slides 
                    for shape in slide.shapes 
                    if hasattr(shape, "text")
                )
                
            elif file_name_lower.endswith(('.csv', '.xlsx', '.xls')):
                if file_name_lower.endswith('.csv'):
                    df = pd.read_csv(io.BytesIO(downloaded_file))
                else:
                    df = pd.read_excel(io.BytesIO(downloaded_file))
                return df.to_string()
                
            else:
                return downloaded_file.decode('utf-8', errors='ignore')
        
        file_content = await asyncio.to_thread(process_file_sync)
        
        if not file_content.strip():
            await send_or_edit_message(
                message.chat.id, 
                f"⚠️ File <code>{escape_html(message.document.file_name)}</code> kosong atau tidak dapat dibaca.", 
                placeholder, 
                parse_mode='HTML'
            )
            return
            
        await context_manager.add_file_to_session(
            message.from_user.id, message.document.file_name, file_content
        )
        await send_or_edit_message(
            message.chat.id, 
            f"✅ Konten dari <code>{escape_html(message.document.file_name)}</code> telah ditambahkan ke sesi.", 
            placeholder, 
            parse_mode="HTML"
        )
        
    except Exception as e:
        logging.error(f"Gagal memproses file {message.document.file_name}: {e}", exc_info=True)
        await send_or_edit_message(message.chat.id, "❌ Gagal memproses file.", placeholder)

async def handle_url_message(message, url):
    placeholder = await bot.reply_to(message, f"🤖🔗 <i>Menganalisis tautan...</i>", parse_mode='HTML')
    content, title = await extract_content_from_url(url)
    
    if not content:
        await send_or_edit_message(
            message.chat.id, "⚠️ Tidak dapat mengambil konten dari tautan.", placeholder
        )
        return
        
    await context_manager.add_file_to_session(
        message.from_user.id, f"URL: {title or url}", content
    )
    await send_or_edit_message(
        message.chat.id, 
        f"✅ Konten dari <b>{escape_html(title)}</b> telah ditambahkan ke sesi.", 
        placeholder, 
        parse_mode="HTML"
    )

async def _handle_menu_action(call, action):
    chat_id = call.message.chat.id
    placeholder = call.message 
    
    if action == "tutup":
        await send_or_edit_message(
            chat_id, "✅ Menu ditutup", placeholder, reply_markup=None, parse_mode="HTML"
        )
        return
        
    if action == "kembali":
        markup = _get_main_menu_markup()
        await send_or_edit_message(
            chat_id, 
            "<b>⚙️ Menu Fitur GHOST</b>\n\nPilih fitur di bawah ini:", 
            placeholder, 
            parse_mode="HTML", 
            reply_markup=markup
        )
        return
        
    if action == "sesi":
        markup = _get_session_menu_markup()
        await send_or_edit_message(
            chat_id, 
            "<b>📁 Menu Sesi</b>\n\nAtur sesi berdasarkan topik/file.", 
            placeholder, 
            parse_mode="HTML", 
            reply_markup=markup
        )
        
    elif action == "reset":
        markup = types.InlineKeyboardMarkup(row_width=2).add(
            types.InlineKeyboardButton("🗑️ Hapus", callback_data="reset_confirm"),
            types.InlineKeyboardButton("⬅️ Kembali", callback_data="menu_kembali")
        )
        await send_or_edit_message(
            chat_id, 
            "<b>⚠️ Hapus Semua</b>\n\nIni akan menghapus semua riwayat & sesi. Data tak bisa dikembalikan. Lanjut?", 
            placeholder, 
            parse_mode="HTML", 
            reply_markup=markup
        )
        
    elif action == "help":
        help_text = (
            "<b>📖 Bantuan & Fitur</b>\n\n"
            "<b>🔹 Interaksi Dasar:</b>\n"
            "• <b>Tanya Apapun:</b> Kirim pesan teks atau suara.\n"
            "• <b>Analisis Gambar:</b> Kirim gambar dan ajukan pertanyaan.\n"
            "• <b>Diskusi File/Tautan:</b> Kirim file atau URL untuk dibahas.\n\n"
            "<b>🎨 Membuat Gambar:</b>\n"
            "Minta saya untuk membuat gambar, contoh: <i>'buatkan gambar astronot di pantai'</i>.\n\n"
            "<b>🔧 Menu Lanjutan:</b>\n"
            "• <b>Kelola Sesi:</b> Lihat konteks aktif atau akhiri sesi.\n"
            "• <b>Hapus Riwayat:</b> Mulai dari awal."
        )
        markup = _get_back_to_main_menu_markup()
        await send_or_edit_message(
            chat_id, help_text, placeholder, parse_mode="HTML", reply_markup=markup
        )


async def _handle_session_menu_action(call, action):
    chat_id = call.message.chat.id
    placeholder = call.message
    
    if action == "aktif":
        context = await context_manager.get_context(call.from_user.id)
        session_name = context.get('active_session_name')
        files_list = context.get('session_files', {}).keys()
        
        if session_name:
            message_text = f"<b>🟢 Sesi Aktif: '{escape_html(session_name)}'</b>\n\n"
        else:
            message_text = "🔘 Tidak ada sesi aktif."
            
        if session_name and files_list:
            files = "\n".join([f"• <code>{escape_html(name)}</code>" for name in files_list])
            message_text += f"<b>📎 Konteks dalam Sesi:</b>\n{files}"
        elif session_name:
            message_text += "<i>Tidak ada file atau tautan dalam sesi ini.</i>"
            
        await send_or_edit_message(
            chat_id, 
            message_text, 
            placeholder, 
            parse_mode="HTML", 
            reply_markup=_get_back_to_session_menu_markup()
        )
        
    elif action == "akhir":
        ended_session = await context_manager.end_session(call.from_user.id)
        if ended_session:
            msg = f"⏹️ Sesi '<b>{escape_html(ended_session)}</b>' telah diakhiri."
        else:
            msg = "❌ Tidak ada sesi aktif."
            
        await send_or_edit_message(
            chat_id, msg, placeholder, parse_mode="HTML", reply_markup=None
        )

async def _handle_reset_action(call, action):
    if action == "confirm":
        if await context_manager.reset_context(call.from_user.id):
            await send_or_edit_message(
                call.message.chat.id, "✅ Semua riwayat telah dihapus.", call.message
            )
        else:
            await send_or_edit_message(
                call.message.chat.id, "ℹ️ Tidak ada riwayat untuk dihapus.", call.message
            )

CALLBACK_ROUTER = {
    'stop': (lambda call, cid: stop_requests.update({int(cid): True})),
    'menu': _handle_menu_action,
    'sesi': _handle_session_menu_action,
    'reset': _handle_reset_action
}

@bot.callback_query_handler(func=lambda call: True)
async def callback_query_dispatcher(call):
    await bot.answer_callback_query(call.id)
    
    try:
        prefix, *data_parts = call.data.split('_', 1)
        handler = CALLBACK_ROUTER.get(prefix)
        
        if handler:
            await handler(call, data_parts[0] if data_parts else "")
        else:
            logging.warning(f"Handler tidak ditemukan untuk prefix: {prefix}")
            
    except Exception as e:
        logging.error(f"Error memproses callback: {call.data}, error: {e}", exc_info=True)